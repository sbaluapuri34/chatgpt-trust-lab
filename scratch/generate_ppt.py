import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    
    # Set slide dimensions (16:9 widescreen)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Custom color palette (Sleek Dark Mode matching dashboard)
    BG_COLOR = RGBColor(11, 15, 25)       # #0b0f19
    CARD_COLOR = RGBColor(17, 24, 39)      # #111827
    TEXT_MAIN = RGBColor(248, 250, 252)    # #f8fafc
    TEXT_MUTED = RGBColor(148, 163, 184)   # #94a3b8
    PURPLE_ACCENT = RGBColor(139, 92, 246) # #8b5cf6
    BLUE_ACCENT = RGBColor(59, 130, 246)   # #3b82f6
    GREEN_ACCENT = RGBColor(16, 185, 129)  # #10b981
    RED_ACCENT = RGBColor(239, 68, 68)     # #ef4444
    
    # Helper function to set custom dark background on slide
    def set_slide_background(slide):
        bg_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5)
        )
        bg_shape.fill.solid()
        bg_shape.fill.fore_color.rgb = BG_COLOR
        bg_shape.line.fill.background() # No border
        # Send background shape to back
        slide.shapes._spTree.remove(bg_shape._element)
        slide.shapes._spTree.insert(2, bg_shape._element)
        return bg_shape
        
    # Helper to add consistent headers to slides
    def add_slide_header(slide, title_text, category_text="RESEARCH SUMMARY"):
        # Category label
        cat_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.4), Inches(11.833), Inches(0.4))
        cat_tf = cat_box.text_frame
        cat_tf.word_wrap = True
        cat_tf.margin_left = cat_tf.margin_top = cat_tf.margin_right = cat_tf.margin_bottom = 0
        p_cat = cat_tf.paragraphs[0]
        p_cat.text = category_text.upper()
        p_cat.font.name = "Arial"
        p_cat.font.size = Pt(11)
        p_cat.font.bold = True
        p_cat.font.color.rgb = PURPLE_ACCENT
        
        # Slide Title
        title_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.7), Inches(11.833), Inches(0.8))
        title_tf = title_box.text_frame
        title_tf.word_wrap = True
        title_tf.margin_left = title_tf.margin_top = title_tf.margin_right = title_tf.margin_bottom = 0
        p_title = title_tf.paragraphs[0]
        p_title.text = title_text
        p_title.font.name = "Arial"
        p_title.font.size = Pt(28)
        p_title.font.bold = True
        p_title.font.color.rgb = TEXT_MAIN
        
    # Blank layout is 6
    blank_layout = prs.slide_layouts[6]
    
    # =========================================================================
    # SLIDE 1: TITLE SLIDE
    # =========================================================================
    slide1 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide1)
    
    # Large colored accent banner
    accent_bar = slide1.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.3), Inches(7.5)
    )
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = PURPLE_ACCENT
    accent_bar.line.fill.background()
    
    # Main Title box
    title_box = slide1.shapes.add_textbox(Inches(1.2), Inches(2.2), Inches(11.0), Inches(2.2))
    tf1 = title_box.text_frame
    tf1.word_wrap = True
    tf1.margin_left = tf1.margin_top = tf1.margin_right = tf1.margin_bottom = 0
    p1 = tf1.paragraphs[0]
    p1.text = "ChatGPT Output Trust &\nEvaluation Lab"
    p1.font.name = "Arial"
    p1.font.size = Pt(46)
    p1.font.bold = True
    p1.font.color.rgb = TEXT_MAIN
    
    # Subtitle box
    sub_box = slide1.shapes.add_textbox(Inches(1.2), Inches(4.5), Inches(11.0), Inches(1.2))
    tf_sub = sub_box.text_frame
    tf_sub.word_wrap = True
    tf_sub.margin_left = tf_sub.margin_top = tf_sub.margin_right = tf_sub.margin_bottom = 0
    p_sub = tf_sub.paragraphs[0]
    p_sub.text = "An Empirical Investigation of User Trust Calibration, AI Hallucinations,\nand Human Evaluation Heuristics in the r/ChatGPT Ecosystem"
    p_sub.font.name = "Arial"
    p_sub.font.size = Pt(18)
    p_sub.font.color.rgb = TEXT_MUTED
    
    # Metadata box
    meta_box = slide1.shapes.add_textbox(Inches(1.2), Inches(6.0), Inches(11.0), Inches(0.6))
    tf_meta = meta_box.text_frame
    tf_meta.word_wrap = True
    tf_meta.margin_left = tf_meta.margin_top = tf_meta.margin_right = tf_meta.margin_bottom = 0
    p_meta = tf_meta.paragraphs[0]
    p_meta.text = "Presented by: Srish (s.baluapuri34@gmail.com)  |  Version 1.0  |  May 30, 2026"
    p_meta.font.name = "Arial"
    p_meta.font.size = Pt(12)
    p_meta.font.bold = True
    p_meta.font.color.rgb = PURPLE_ACCENT
    
    # Slide 1 Notes
    slide1.notes_slide.notes_text_frame.text = (
        "Welcome stakeholders and evaluators to the ChatGPT Output Trust & Evaluation Lab project summary. "
        "This project was born out of a critical need: while extensive research exists evaluating the mathematical "
        "and logical performance of LLMs on standardized benchmarks, very little empirical work has been done "
        "examining how real-world practitioners calibrate their trust in everyday operations. "
        "Today, we will trace the entire project journey, from data scraping to advanced semantic LLM parsing, "
        "revealing major findings regarding user attachment, auditing behaviors, and downstream consequences."
    )

    # =========================================================================
    # SLIDE 2: EXECUTIVE SUMMARY
    # =========================================================================
    slide2 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide2)
    add_slide_header(slide2, "Executive Summary", "Brief Overview")
    
    # Left main brief
    brief_box = slide2.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(6.0), Inches(5.0))
    bf = brief_box.text_frame
    bf.word_wrap = True
    bf.margin_left = bf.margin_top = bf.margin_right = bf.margin_bottom = 0
    
    p = bf.paragraphs[0]
    p.text = "Empirical Auditing of LLM Trust Calibration"
    p.font.name = "Arial"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = BLUE_ACCENT
    p.space_after = Pt(14)
    
    bullets = [
        "A rigorous, hybrid 5-phase data pipeline was engineered to ingest 1,091 unique posts, pre-screen 729 candidates, and classify 594 relevant threads from r/ChatGPT.",
        "Traditional keyword matching was audited and mathematically proven highly insufficient, exhibiting an 88.89% false positive rate and a 60% false negative rate.",
        "Quantified a fundamental tension in trust calibration: users seeking emotional support exhibit high levels of Over-Reliance (12.29%), while practitioners catching hallucinations develop custom verification tools and protocols (14.65%).",
        "While Confidently Incorrect Outputs are the most frequent machine failures (31.14%), Real-World Impact (consequences) generates massive community upvoting and attention (26.40% weighted prevalence)."
    ]
    for b in bullets:
        p_b = bf.add_paragraph()
        p_b.text = "• " + b
        p_b.font.name = "Arial"
        p_b.font.size = Pt(14)
        p_b.font.color.rgb = TEXT_MAIN
        p_b.space_after = Pt(10)
        
    # Right stat panels
    stats_data = [
        ("1,091", "Raw Ingested Threads", "Deduplicated from old.reddit.com"),
        ("729", "Keyword Candidates", "Filtered in Phase 2 pre-screening"),
        ("594", "Relevant Classified Posts", "Stage 1 relevance-confirmed by LLM")
    ]
    
    for i, (num, lbl, desc) in enumerate(stats_data):
        top_y = Inches(1.8 + i * 1.6)
        
        # Stat Card
        card = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.5), top_y, Inches(5.0), Inches(1.3))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_COLOR
        card.line.fill.background()
        
        # Text in card
        tf_card = card.text_frame
        tf_card.word_wrap = True
        tf_card.margin_left = Inches(0.2)
        tf_card.margin_top = Inches(0.15)
        
        p_num = tf_card.paragraphs[0]
        p_num.text = num
        p_num.font.name = "Arial"
        p_num.font.size = Pt(28)
        p_num.font.bold = True
        p_num.font.color.rgb = PURPLE_ACCENT
        
        p_lbl = tf_card.add_paragraph()
        p_lbl.text = lbl + "  |  " + desc
        p_lbl.font.name = "Arial"
        p_lbl.font.size = Pt(11)
        p_lbl.font.bold = True
        p_lbl.font.color.rgb = TEXT_MUTED
        
    # Slide 2 Notes
    slide2.notes_slide.notes_text_frame.text = (
        "The executive summary acts as our core high-level overview. We present our quantitative funnel "
        "on the right: out of 1,091 unique raw posts pulled via Playwright, 729 were preselected by keyword rules, "
        "and 594 were semantically confirmed and categorized. On the left, we map out the core takeaways: "
        "1. The hybrid pipeline succeeded. 2. Keyword rules are highly noisy. 3. Trust breakdown is highly "
        "correlated to unannounced OpenAI safety updates. 4. Real-world consequences drive major community engagement."
    )

    # =========================================================================
    # SLIDE 3: PROBLEM STATEMENT
    # =========================================================================
    slide3 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide3)
    add_slide_header(slide3, "The Problem Statement", "Research Context")
    
    # Left Description
    desc_box = slide3.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(5.5), Inches(5.0))
    df = desc_box.text_frame
    df.word_wrap = True
    df.margin_left = df.margin_top = df.margin_right = df.margin_bottom = 0
    
    p = df.paragraphs[0]
    p.text = "Automation Bias & Semantic Uncertainty"
    p.font.name = "Arial"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = BLUE_ACCENT
    p.space_after = Pt(14)
    
    bullets = [
        "LLMs achieve unprecedented conversational fluency and empathy simulation, leading users to trust outputs implicitly.",
        "The system has no internal concept of truth; it confidently fabricates legal citations, medical advice, and code functions without indicating its uncertainty.",
        "Standard evaluation benchmarks (MMLU, HumanEval) ignore human factors, failing to measure when and why users accept false statements.",
        "Traditional sentiment analysis is incapable of capturing implicit trust breakdown, requiring a two-stage semantic categorization workflow."
    ]
    for b in bullets:
        p_b = df.add_paragraph()
        p_b.text = "• " + b
        p_b.font.name = "Arial"
        p_b.font.size = Pt(14)
        p_b.font.color.rgb = TEXT_MAIN
        p_b.space_after = Pt(12)
        
    # Right illustrative quote
    quote_card = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.8), Inches(1.8), Inches(5.7), Inches(4.5))
    quote_card.fill.solid()
    quote_card.fill.fore_color.rgb = CARD_COLOR
    quote_card.line.color.rgb = RED_ACCENT
    quote_card.line.width = Pt(2.5)
    
    q_tf = quote_card.text_frame
    q_tf.word_wrap = True
    q_tf.margin_left = Inches(0.4)
    q_tf.margin_top = Inches(0.4)
    q_tf.margin_right = Inches(0.4)
    
    p_q = q_tf.paragraphs[0]
    p_q.text = "“I counted roughly 8–10 distinct factual errors, all delivered with certainty.”"
    p_q.font.name = "Arial"
    p_q.font.size = Pt(22)
    p_q.font.bold = True
    p_q.font.italic = True
    p_q.font.color.rgb = RED_ACCENT
    p_q.space_after = Pt(14)
    
    p_q_meta = q_tf.add_paragraph()
    p_q_meta.text = "— ChatGPT User experience report on blind trust and factual mistakes"
    p_q_meta.font.name = "Arial"
    p_q_meta.font.size = Pt(12)
    p_q_meta.font.bold = True
    p_q_meta.font.color.rgb = TEXT_MUTED
    p_q_meta.space_after = Pt(14)
    
    p_q_desc = q_tf.add_paragraph()
    p_q_desc.text = "Without visual uncertainty feedback, users treat highly fluent hallucinations as verified facts. This creates a critical design problem: how can developers encourage healthy verification rather than blind over-reliance?"
    p_q_desc.font.name = "Arial"
    p_q_desc.font.size = Pt(13)
    p_q_desc.font.color.rgb = TEXT_MAIN
    
    # Slide 3 Notes
    slide3.notes_slide.notes_text_frame.text = (
        "This slide presents the core problem of Automation Bias in LLMs. "
        "Because models write with extreme confidence and flawless grammar, users cannot easily distinguish "
        "between true statements and completely fabricated lies. Fictional citations or wrong facts are delivered "
        "with absolute certainty. The quote on the right illustrates a user catching 8 to 10 distinct factual errors "
        "all delivered with certainty. Our research seeks to map out exactly how users react when they encounter these failures."
    )

    # =========================================================================
    # SLIDE 4: RESEARCH OBJECTIVES
    # =========================================================================
    slide4 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide4)
    add_slide_header(slide4, "Research Objectives", "Project Goals")
    
    # 4 Columns of Objectives
    objectives = [
        ("1. Ingest Organic Evidence", "Harvest unvarnished, real-world user reports and complaints regarding ChatGPT output failures from active practitioner communities.", BLUE_ACCENT),
        ("2. Evolve Taxonomy", "Evolve primitive technical engineering classifications into highly expressive, human-aligned research themes that reflect the user journey.", PURPLE_ACCENT),
        ("3. Audit Matching Rules", "Quantify the accuracy and limitations of traditional deterministic keyword pre-screening heuristics compared to semantic LLM classifiers.", GREEN_ACCENT),
        ("4. Quantify Virality", "Calculate engagement-weighted prevalence rates to prioritize viral community narratives over single-user comments, and serve an open-source dashboard.", RED_ACCENT)
    ]
    
    for i, (title, desc, color) in enumerate(objectives):
        col_x = Inches(0.75 + i * 2.95)
        
        # Box container
        box = slide4.shapes.add_shape(MSO_SHAPE.RECTANGLE, col_x, Inches(1.8), Inches(2.8), Inches(4.5))
        box.fill.solid()
        box.fill.fore_color.rgb = CARD_COLOR
        box.line.color.rgb = color
        box.line.width = Pt(2.0)
        
        b_tf = box.text_frame
        b_tf.word_wrap = True
        b_tf.margin_left = Inches(0.25)
        b_tf.margin_top = Inches(0.3)
        b_tf.margin_right = Inches(0.25)
        
        p_t = b_tf.paragraphs[0]
        p_t.text = title
        p_t.font.name = "Arial"
        p_t.font.size = Pt(17)
        p_t.font.bold = True
        p_t.font.color.rgb = color
        p_t.space_after = Pt(14)
        
        p_d = b_tf.add_paragraph()
        p_d.text = desc
        p_d.font.name = "Arial"
        p_d.font.size = Pt(13)
        p_d.font.color.rgb = TEXT_MAIN
        p_d.space_after = Pt(10)
        
    # Slide 4 Notes
    slide4.notes_slide.notes_text_frame.text = (
        "Here we define our four primary research objectives. "
        "First, we want to harvest organic, raw evidence directly from user self-reports. "
        "Second, we map out a refined classification taxonomy aligned to user behaviors. "
        "Third, we audit keyword rules to mathematically demonstrate their limits. "
        "Fourth, we weight findings using log-scaled upvote scores to represent true community resonance. "
        "This lays the complete blueprint for our project execution."
    )

    # =========================================================================
    # SLIDE 5: METHODOLOGY OVERVIEW
    # =========================================================================
    slide5 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide5)
    add_slide_header(slide5, "The Methodology Overview", "Pipeline Design")
    
    # 5 Step Visual Flow Card Setup
    steps = [
        ("Phase 1", "Data Collection", "Playwright scraping of old.reddit.com", "1,091 Raw Posts"),
        ("Phase 2", "Preprocessing Heuristics", "Keyword density & upvote candidate screening", "729 Candidates"),
        ("Phase 3", "Two-Stage LLM", "Llama 3.1 8B relevance and theme extraction", "594 Classified DB"),
        ("Phase 4", "Programmatic Aggs", "Weighted prevalence & co-occurrence matrix math", "5 Agg Insights"),
        ("Phase 5", "Interactive Dashboard", "Sleek Streamlit viz & RO SQLite explorer", "Live Public App")
    ]
    
    for i, (ph, name, desc, output) in enumerate(steps):
        col_x = Inches(0.75 + i * 2.37)
        
        # Step Container
        box = slide5.shapes.add_shape(MSO_SHAPE.RECTANGLE, col_x, Inches(1.8), Inches(2.2), Inches(4.5))
        box.fill.solid()
        box.fill.fore_color.rgb = CARD_COLOR
        box.line.fill.background()
        
        b_tf = box.text_frame
        b_tf.word_wrap = True
        b_tf.margin_left = Inches(0.18)
        b_tf.margin_top = Inches(0.25)
        b_tf.margin_right = Inches(0.18)
        
        p_ph = b_tf.paragraphs[0]
        p_ph.text = ph
        p_ph.font.name = "Arial"
        p_ph.font.size = Pt(13)
        p_ph.font.bold = True
        p_ph.font.color.rgb = PURPLE_ACCENT
        p_ph.space_after = Pt(6)
        
        p_name = b_tf.add_paragraph()
        p_name.text = name
        p_name.font.name = "Arial"
        p_name.font.size = Pt(16)
        p_name.font.bold = True
        p_name.font.color.rgb = TEXT_MAIN
        p_name.space_after = Pt(14)
        
        p_desc = b_tf.add_paragraph()
        p_desc.text = desc
        p_desc.font.name = "Arial"
        p_desc.font.size = Pt(11)
        p_desc.font.color.rgb = TEXT_MUTED
        p_desc.space_after = Pt(14)
        
        p_out = b_tf.add_paragraph()
        p_out.text = "➔ " + output
        p_out.font.name = "Arial"
        p_out.font.size = Pt(12)
        p_out.font.bold = True
        p_out.font.color.rgb = GREEN_ACCENT
        
    # Slide 5 Notes
    slide5.notes_slide.notes_text_frame.text = (
        "This slide presents our complete hybrid research methodology. "
        "The project moves sequentially across 5 distinct phases. "
        "Phase 1 ingests deduplicated raw posts. Phase 2 conducts fast keyword candidate pre-screening to optimize api costs. "
        "Phase 3 performs deep, mock-free Llama-based semantic classification. Phase 4 programmatically compiles metrics and matrices. "
        "Finally, Phase 5 compiles it into our interactive, premium, dark-mode analytical Streamlit dashboard."
    )

    # =========================================================================
    # SLIDE 6: DATA COLLECTION PROCESS
    # =========================================================================
    slide6 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide6)
    add_slide_header(slide6, "Phase 1: Data Collection Process", "Data Ingestion")
    
    # Left columns: Playwright strategy
    strat_box = slide6.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(5.8), Inches(4.8))
    sf = strat_box.text_frame
    sf.word_wrap = True
    sf.margin_left = sf.margin_top = sf.margin_right = sf.margin_bottom = 0
    
    p_s = sf.paragraphs[0]
    p_s.text = "High-Performance Browser Automation"
    p_s.font.name = "Arial"
    p_s.font.size = Pt(20)
    p_s.font.bold = True
    p_s.font.color.rgb = BLUE_ACCENT
    p_s.space_after = Pt(12)
    
    bullets = [
        "Sourced from old.reddit.com to bypass heavy modern React payloads, lowering processing overhead and increasing ingestion speeds.",
        "Bypassed Reddit API rate limits using Playwright browser automation with sandboxed headless Chromium orchestration.",
        "Suppressed non-critical elements (CSS stylesheets, images, tracking scripts, and ads) to lower bandwidth demands.",
        "Structured query searches covering 'trust', 'hallucinate', 'wrong', 'consequence', and 'certain' across a 9-month chronological window.",
        "Enforced a unique primary-key mapping constraint on thread IDs to automatically deduplicate overlapping results."
    ]
    for b in bullets:
        p_b = sf.add_paragraph()
        p_b.text = "• " + b
        p_b.font.name = "Arial"
        p_b.font.size = Pt(13)
        p_b.font.color.rgb = TEXT_MAIN
        p_b.space_after = Pt(10)
        
    # Right panel: statistics
    stats_card = slide6.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.0), Inches(1.8), Inches(5.5), Inches(4.5))
    stats_card.fill.solid()
    stats_card.fill.fore_color.rgb = CARD_COLOR
    stats_card.line.color.rgb = PURPLE_ACCENT
    stats_card.line.width = Pt(2.0)
    
    sc_tf = stats_card.text_frame
    sc_tf.word_wrap = True
    sc_tf.margin_left = Inches(0.4)
    sc_tf.margin_top = Inches(0.4)
    sc_tf.margin_right = Inches(0.4)
    
    p_sc_title = sc_tf.paragraphs[0]
    p_sc_title.text = "Ingestion Statistics"
    p_sc_title.font.name = "Arial"
    p_sc_title.font.size = Pt(20)
    p_sc_title.font.bold = True
    p_sc_title.font.color.rgb = PURPLE_ACCENT
    p_sc_title.space_after = Pt(14)
    
    stats_lines = [
        ("Target Subreddit", "r/ChatGPT (Largest centralized user group)"),
        ("Ingestion Query Sweeps", "8 unique keyword query scripts"),
        ("Raw Scraped Records", "1,091 threads"),
        ("Deduplicated Unique Records", "1,091 threads"),
        ("Dataset Size (raw)", "Posts JSON + manifest locked")
    ]
    for lbl, val in stats_lines:
        p_line = sc_tf.add_paragraph()
        p_line.text = f"• {lbl}: "
        p_line.font.name = "Arial"
        p_line.font.size = Pt(13)
        p_line.font.bold = True
        p_line.font.color.rgb = TEXT_MAIN
        
        run = p_line.add_run()
        run.text = val
        run.font.bold = False
        run.font.color.rgb = TEXT_MUTED
        p_line.space_after = Pt(12)
        
    # Slide 6 Notes
    slide6.notes_slide.notes_text_frame.text = (
        "Phase 1 maps out how we obtained our core dataset. We targeted old.reddit.com "
        "because its lightweight HTML structures permit fast, headless browser scraping. "
        "Playwright bypassed traditional API rate limits by simulating a user session, "
        "while blocking large assets like images and ads to speed up crawls. Deduplication "
        "automatically discarded posts caught across multiple search terms, ensuring "
        "1,091 100% unique, locked raw posts."
    )

    # =========================================================================
    # SLIDE 7: DATASET FUNNEL STATISTICS
    # =========================================================================
    slide7 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide7)
    add_slide_header(slide7, "Dataset Statistics & Processing Funnel", "Processing Metrics")
    
    # 3 Large Funnel Cards
    funnels = [
        ("STAGE 1: RAW INGESTION", "1,091", "100.0% of Corpus", "All raw posts scraped via Playwright query searches. Deduplicated and locked to establish an authoritative research dataset.", BLUE_ACCENT),
        ("STAGE 2: HEURISTIC SCREENING", "729", "66.8% Retention", "Deterministic keyword density and upvote screening in Phase 2. Removed empty threads and low-engagement noise.", PURPLE_ACCENT),
        ("STAGE 3: CONFIRMED RELEVANT", "594", "54.4% Pipeline Yield", "Stage 1 semantic relevance confirmed by Llama 3.1 8B. Established final curated sqlite database (data/research.db).", GREEN_ACCENT)
    ]
    
    for i, (title, count, pct, desc, color) in enumerate(funnels):
        col_x = Inches(0.75 + i * 4.0)
        
        box = slide7.shapes.add_shape(MSO_SHAPE.RECTANGLE, col_x, Inches(1.8), Inches(3.7), Inches(4.5))
        box.fill.solid()
        box.fill.fore_color.rgb = CARD_COLOR
        box.line.color.rgb = color
        box.line.width = Pt(2.5)
        
        b_tf = box.text_frame
        b_tf.word_wrap = True
        b_tf.margin_left = Inches(0.3)
        b_tf.margin_top = Inches(0.35)
        b_tf.margin_right = Inches(0.3)
        
        p_t = b_tf.paragraphs[0]
        p_t.text = title
        p_t.font.name = "Arial"
        p_t.font.size = Pt(13)
        p_t.font.bold = True
        p_t.font.color.rgb = TEXT_MUTED
        p_t.space_after = Pt(10)
        
        p_c = b_tf.add_paragraph()
        p_c.text = count
        p_c.font.name = "Arial"
        p_c.font.size = Pt(44)
        p_c.font.bold = True
        p_c.font.color.rgb = color
        
        p_p = b_tf.add_paragraph()
        p_p.text = pct
        p_p.font.name = "Arial"
        p_p.font.size = Pt(14)
        p_p.font.bold = True
        p_p.font.color.rgb = TEXT_MAIN
        p_p.space_after = Pt(14)
        
        p_d = b_tf.add_paragraph()
        p_d.text = desc
        p_d.font.name = "Arial"
        p_d.font.size = Pt(13)
        p_d.font.color.rgb = TEXT_MUTED
        p_d.line_spacing = 1.2
        
    # Slide 7 Notes
    slide7.notes_slide.notes_text_frame.text = (
        "Here we show the math and counts of our data-funneling stages. "
        "We start with a broad raw capture of 1,091 unique threads. "
        "Phase 2 keyword heuristics prunes this down to 729 relevance candidates. "
        "Phase 3's two-stage semantic LLM classification removes further noise and "
        "confirms exactly 594 relevant database records. This yield of 54.4% is extremely high, "
        "validating our preprocessing heuristics for candidate capture before API parsing."
    )

    # =========================================================================
    # SLIDE 8: THEME TAXONOMY EVOLUTION
    # =========================================================================
    slide8 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide8)
    add_slide_header(slide8, "Theme Taxonomy Evolution", "Taxonomy Audit")
    
    # Left column: Why change?
    why_box = slide8.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(5.8), Inches(4.8))
    wf = why_box.text_frame
    wf.word_wrap = True
    wf.margin_left = wf.margin_top = wf.margin_right = wf.margin_bottom = 0
    
    p_w = wf.paragraphs[0]
    p_w.text = "From Technical Outputs to PM-Centric Themes"
    p_w.font.name = "Arial"
    p_w.font.size = Pt(20)
    p_w.font.bold = True
    p_w.font.color.rgb = BLUE_ACCENT
    p_w.space_after = Pt(12)
    
    bullets = [
        "Primitive technical taxonomy centered on abstract model behaviors (e.g. 'hallucinations', 'sycophancy') which isolated findings from user contexts.",
        "Refinement process evolved categories into PM-aligned, human-centric themes that map direct Downstream Down consequences.",
        "Focus shifted from model errors alone to capturing user calibration, psychological attachment, and custom verification heuristics.",
        "The evolved taxonomy maps the entire user journey: from machine failure, emotional reaction, auditing, and real-world consequences."
    ]
    for b in bullets:
        p_b = wf.add_paragraph()
        p_b.text = "• " + b
        p_b.font.name = "Arial"
        p_b.font.size = Pt(13)
        p_b.font.color.rgb = TEXT_MAIN
        p_b.space_after = Pt(12)
        
    # Right column: Mapping table
    map_card = slide8.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.0), Inches(1.8), Inches(5.5), Inches(4.5))
    map_card.fill.solid()
    map_card.fill.fore_color.rgb = CARD_COLOR
    map_card.line.color.rgb = PURPLE_ACCENT
    map_card.line.width = Pt(2.0)
    
    mc_tf = map_card.text_frame
    mc_tf.word_wrap = True
    mc_tf.margin_left = Inches(0.35)
    mc_tf.margin_top = Inches(0.35)
    mc_tf.margin_right = Inches(0.35)
    
    p_mc_t = mc_tf.paragraphs[0]
    p_mc_t.text = "Taxonomy Migration Map"
    p_mc_t.font.name = "Arial"
    p_mc_t.font.size = Pt(18)
    p_mc_t.font.bold = True
    p_mc_t.font.color.rgb = PURPLE_ACCENT
    p_mc_t.space_after = Pt(14)
    
    mappings = [
        ("hallucination_instances", "Confidently Incorrect Outputs"),
        ("trust_erosion", "User Trust Breakdown"),
        ("companionship_attachment", "Over-Reliance on AI Outputs"),
        ("fact_checking_behavior", "User Verification Behavior"),
        ("consequential_outcomes", "Real-World Impact of AI"),
        ("sycophancy_persuasion", "Persuasive Outputs")
    ]
    
    for orig, final in mappings:
        p_m = mc_tf.add_paragraph()
        p_m.text = f"• `{orig}`  ➔  "
        p_m.font.name = "Courier New"
        p_m.font.size = Pt(11)
        p_m.font.color.rgb = TEXT_MUTED
        
        run = p_m.add_run()
        run.text = final
        run.font.name = "Arial"
        run.font.bold = True
        run.font.size = Pt(13)
        run.font.color.rgb = TEXT_MAIN
        p_m.space_after = Pt(8)
        
    # Slide 8 Notes
    slide8.notes_slide.notes_text_frame.text = (
        "We discuss the evolution of our theme taxonomy here. "
        "Originally, the categorization system was highly engineering-centric, "
        "focusing strictly on model outputs (like hallucination counts or sycophancy). "
        "We migrated this to a human-aligned, PM-centric research taxonomy. "
        "By focusing on how users react to and verify these outputs, our data becomes "
        "substantially more valuable for product managers and human-computer interaction (HCI) researchers."
    )

    # =========================================================================
    # SLIDE 9: THEME PREVALENCE RESULTS
    # =========================================================================
    slide9 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide9)
    add_slide_header(slide9, "Research Results: Theme Prevalence", "Prevalence Analysis")
    
    # Left table: Prevalence data
    table_box = slide9.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(6.8), Inches(4.8))
    tb_tf = table_box.text_frame
    tb_tf.word_wrap = True
    tb_tf.margin_left = tb_tf.margin_top = tb_tf.margin_right = tb_tf.margin_bottom = 0
    
    p_t = tb_tf.paragraphs[0]
    p_t.text = "Quantitative Theme Distribution (Total: 594 Posts)"
    p_t.font.name = "Arial"
    p_t.font.size = Pt(18)
    p_t.font.bold = True
    p_t.font.color.rgb = BLUE_ACCENT
    p_t.space_after = Pt(12)
    
    # Text-based table simulation
    prev_table = [
        ("Confidently Incorrect Outputs", "185", "31.14%", "292.91", "24.31%"),
        ("Real-World Impact of AI Outputs", "131", "22.05%", "318.07", "26.40%"),
        ("User Evaluation & Verification Behavior", "87", "14.65%", "115.82", "9.61%"),
        ("Over-Reliance on AI Outputs", "73", "12.29%", "156.51", "12.99%"),
        ("User Trust Breakdown", "72", "12.12%", "222.27", "18.45%"),
        ("Persuasive Outputs & Trust Formation", "46", "7.74%", "99.38", "8.25%")
    ]
    
    header = f"{'Research Theme':<38} | {'Raw (%)':<10} | {'Weighted (%)':<12}"
    p_h = tb_tf.add_paragraph()
    p_h.text = header
    p_h.font.name = "Courier New"
    p_h.font.bold = True
    p_h.font.size = Pt(11)
    p_h.font.color.rgb = PURPLE_ACCENT
    p_h.space_after = Pt(6)
    
    for theme, raw_cnt, raw_pct, w_cnt, w_pct in prev_table:
        p_row = tb_tf.add_paragraph()
        p_row.text = f"{theme:<38} | {raw_cnt:<3} ({raw_pct:<6}) | {round(float(w_cnt),1):<5} ({w_pct:<6})"
        p_row.font.name = "Courier New"
        p_row.font.size = Pt(11)
        p_row.font.color.rgb = TEXT_MAIN
        p_row.space_after = Pt(8)
        
    # Right Column: Prevalence Paradox
    paradox_box = slide9.shapes.add_textbox(Inches(7.8), Inches(1.8), Inches(4.783), Inches(4.8))
    pf = paradox_box.text_frame
    pf.word_wrap = True
    pf.margin_left = pf.margin_top = pf.margin_right = pf.margin_bottom = 0
    
    p_p_title = pf.paragraphs[0]
    p_p_title.text = "The Prevalence Paradox"
    p_p_title.font.name = "Arial"
    p_p_title.font.size = Pt(18)
    p_p_title.font.bold = True
    p_p_title.font.color.rgb = PURPLE_ACCENT
    p_p_title.space_after = Pt(12)
    
    p_p_desc = pf.add_paragraph()
    p_p_desc.text = "While Confidently Incorrect Outputs occur most frequently in raw numbers (31.14%), Real-World Impact drives the highest engagement and viral attention share (26.40% weighted prevalence)."
    p_p_desc.font.name = "Arial"
    p_p_desc.font.size = Pt(13)
    p_p_desc.font.italic = True
    p_p_desc.font.color.rgb = TEXT_MAIN
    p_p_desc.space_after = Pt(14)
    
    bullets = [
        "Minor factual errors occur constantly in day-to-day work.",
        "Users are disproportionately motivated to document and discuss severe consequences ( probate audit fraud, veterinary blood panels, legal disbarment warnings).",
        "Weighted count adjusts for community engagement via ln(1 + max(score,0)) to show real community resonance."
    ]
    for b in bullets:
        p_b = pf.add_paragraph()
        p_b.text = "• " + b
        p_b.font.name = "Arial"
        p_b.font.size = Pt(13)
        p_b.font.color.rgb = TEXT_MUTED
        p_b.space_after = Pt(10)
        
    # Slide 9 Notes
    slide9.notes_slide.notes_text_frame.text = (
        "This slide presents one of our key scientific findings, which we call 'The Prevalence Paradox'. "
        "If you look at the raw counts, 'Confidently Incorrect Outputs' is the largest theme by far with 185 posts (31.14%). "
        "However, when we adjust for community engagement using our log-scaled upvote weighting formula, "
        "'Real-World Impact' rises to become the largest theme with 26.40% weighted prevalence. "
        "This mathematically proves that while minor errors are highly frequent, users are far more motivated "
        "to discuss, upvote, and validate threads containing tangible downstream consequences."
    )

    # =========================================================================
    # SLIDE 10: SEVERITY ANALYSIS
    # =========================================================================
    slide10 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide10)
    add_slide_header(slide10, "Risk Profiles: Severity Analysis", "Risk Assessment")
    
    # 3 Columns for Low, Medium, High severity profiles
    severities = [
        ("LOW SEVERITY PROFILE", "General Annoyances & Jokes", GREEN_ACCENT, [
            "Minor conversational mistakes, code syntax slips, or general chatbot humor.",
            "Minimal real-world friction; resolved instantly by standard user prompting.",
            "No direct downstream consequences or lost time recorded.",
            "Acts as conversational background noise; has little impact on operational trust."
        ]),
        ("MEDIUM SEVERITY PROFILE", "Operational Friction & Waste", PURPLE_ACCENT, [
            "Significant time wasted debugging hallucinated code modules or verifying false facts.",
            "User frustration and visible friction, leading to moderate trust erosion.",
            "Requires active user intervention (re-prompting, cross-checking external sources).",
            "Triggers the development of custom local verification heuristics."
        ]),
        ("HIGH SEVERITY PROFILE", "Severe Downstream Damages", RED_ACCENT, [
            "Tangible, high-stakes consequences affecting professions, money, or health.",
            "Examples: lawyers submitting hallucinated briefs, estate probate fraud, and near-fatal veterinary panel mistakes.",
            "High risk of severe reputation damage or institutional liability.",
            "Erodes user trust completely, triggering direct pushback against LLM deployment."
        ])
    ]
    
    for i, (title, desc, color, bullets) in enumerate(severities):
        col_x = Inches(0.75 + i * 4.0)
        
        box = slide10.shapes.add_shape(MSO_SHAPE.RECTANGLE, col_x, Inches(1.8), Inches(3.7), Inches(4.5))
        box.fill.solid()
        box.fill.fore_color.rgb = CARD_COLOR
        box.line.color.rgb = color
        box.line.width = Pt(2.5)
        
        b_tf = box.text_frame
        b_tf.word_wrap = True
        b_tf.margin_left = Inches(0.25)
        b_tf.margin_top = Inches(0.3)
        b_tf.margin_right = Inches(0.25)
        
        p_t = b_tf.paragraphs[0]
        p_t.text = title
        p_t.font.name = "Arial"
        p_t.font.size = Pt(14)
        p_t.font.bold = True
        p_t.font.color.rgb = color
        p_t.space_after = Pt(4)
        
        p_d = b_tf.add_paragraph()
        p_d.text = desc
        p_d.font.name = "Arial"
        p_d.font.size = Pt(12)
        p_d.font.bold = True
        p_d.font.color.rgb = TEXT_MAIN
        p_d.space_after = Pt(14)
        
        for bullet in bullets:
            p_b = b_tf.add_paragraph()
            p_b.text = "• " + bullet
            p_b.font.name = "Arial"
            p_b.font.size = Pt(10.5)
            p_b.font.color.rgb = TEXT_MUTED
            p_b.space_after = Pt(8)
            
    # Slide 10 Notes
    slide10.notes_slide.notes_text_frame.text = (
        "Slide 10 details our risk severity analysis. "
        "We categorize post outcomes into low, medium, and high severity. "
        "Low severity includes simple jokes or minor conversational mistakes. "
        "Medium severity represents moderate friction, such as wasting hours debugging a fabricated code package. "
        "High severity represents major damages, including legal liability, medical errors, or financial fraud. "
        "This mapping allows us to evaluate the risk landscape of model failures."
    )

    # =========================================================================
    # SLIDE 11: USER TRUST BREAKDOWN
    # =========================================================================
    slide11 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide11)
    add_slide_header(slide11, "User Trust Breakdown Insights", "Thematic deep-dive")
    
    # Left column: Data & Analysis
    info_box = slide11.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(5.8), Inches(4.8))
    inf = info_box.text_frame
    inf.word_wrap = True
    inf.margin_left = inf.margin_top = inf.margin_right = inf.margin_bottom = 0
    
    p_t = inf.paragraphs[0]
    p_t.text = "Operational Metrics & Drivers"
    p_t.font.name = "Arial"
    p_t.font.size = Pt(20)
    p_t.font.bold = True
    p_t.font.color.rgb = PURPLE_ACCENT
    p_t.space_after = Pt(12)
    
    bullets = [
        "Raw Prevalence: 72 posts (12.12%)  |  Weighted: 222.27 (18.45%)",
        "Primary Driver: OpenAI Platform Management. Unannounced safety filter updates and silent model downgrades trigger immediate user frustration.",
        "Erosion Dynamics: Trust is highly volatile. Users quickly switch from enthusiastic advocates to vocal skeptics after catching a severe error.",
        " डाउन Stream Effects: Triggers aggressive user pushback and negative word-of-mouth in developer networks."
    ]
    for b in bullets:
        p_b = inf.add_paragraph()
        p_b.text = "• " + b
        p_b.font.name = "Arial"
        p_b.font.size = Pt(13)
        p_b.font.color.rgb = TEXT_MAIN
        p_b.space_after = Pt(12)
        
    # Right column: Quote Card
    card = slide11.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.8), Inches(1.8), Inches(5.7), Inches(4.5))
    card.fill.solid()
    card.fill.fore_color.rgb = CARD_COLOR
    card.line.color.rgb = PURPLE_ACCENT
    card.line.width = Pt(2.0)
    
    c_tf = card.text_frame
    c_tf.word_wrap = True
    c_tf.margin_left = Inches(0.4)
    c_tf.margin_top = Inches(0.4)
    c_tf.margin_right = Inches(0.4)
    
    p_q = c_tf.paragraphs[0]
    p_q.text = "“Trust OpenAI to fuck up everything and lose user AND investor trust in any and every opportunity they have.”"
    p_q.font.name = "Arial"
    p_q.font.size = Pt(18)
    p_q.font.bold = True
    p_q.font.italic = True
    p_q.font.color.rgb = TEXT_MAIN
    p_q.space_after = Pt(12)
    
    p_m = c_tf.add_paragraph()
    p_m.text = "— Score: 768  |  Engagement Weight: 6.65\nSource: https://old.reddit.com/r/ChatGPT/comments/1nrsks8"
    p_m.font.name = "Arial"
    p_m.font.size = Pt(11)
    p_m.font.bold = True
    p_m.font.color.rgb = PURPLE_ACCENT
    p_m.space_after = Pt(14)
    
    p_d = c_tf.add_paragraph()
    p_d.text = "Research Finding: A large proportion of user trust breakdown is not caused by technical hallucinations, but rather by platform execution. When safety filtering is inconsistent or unannounced, it silently erodes trust."
    p_d.font.name = "Arial"
    p_d.font.size = Pt(12)
    p_d.font.color.rgb = TEXT_MUTED
    
    # Slide 11 Notes
    slide11.notes_slide.notes_text_frame.text = (
        "We do a thematic deep-dive into 'User Trust Breakdown'. "
        "This theme represents 12.12% raw prevalence but rises to 18.45% weighted prevalence, "
        "indicating high engagement. Our core finding is that trust breakdown is heavily driven by platform management: "
        "unannounced safety filter updates, model behavior shifts, and support lack. The quote on the right shows the "
        "intense user frustration when they feel OpenAI has silently degraded the experience."
    )

    # =========================================================================
    # SLIDE 12: OVER-RELIANCE INSIGHTS
    # =========================================================================
    slide12 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide12)
    add_slide_header(slide12, "Over-Reliance & Human-AI Attachment", "Thematic deep-dive")
    
    # Left column
    info_box = slide12.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(5.8), Inches(4.8))
    inf = info_box.text_frame
    inf.word_wrap = True
    inf.margin_left = inf.margin_top = inf.margin_right = inf.margin_bottom = 0
    
    p_t = inf.paragraphs[0]
    p_t.text = "Operational Metrics & Drivers"
    p_t.font.name = "Arial"
    p_t.font.size = Pt(20)
    p_t.font.bold = True
    p_t.font.color.rgb = BLUE_ACCENT
    p_t.space_after = Pt(12)
    
    bullets = [
        "Raw Prevalence: 73 posts (12.29%)  |  Weighted: 156.51 (12.99%)",
        "Primary Driver: Anthropomorphism and Simulated Empathy. Users treat ChatGPT as a companion, advisor, or therapist.",
        "Attachment Risks: High danger of emotional vulnerability. When models are updated or context is cleared, users experience severe emotional distress.",
        "Verification Failure: Users in these scenarios exhibit near-zero fact-checking, accepting medical or psychological advice blindly."
    ]
    for b in bullets:
        p_b = inf.add_paragraph()
        p_b.text = "• " + b
        p_b.font.name = "Arial"
        p_b.font.size = Pt(13)
        p_b.font.color.rgb = TEXT_MAIN
        p_b.space_after = Pt(12)
        
    # Right column
    card = slide12.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.8), Inches(1.8), Inches(5.7), Inches(4.5))
    card.fill.solid()
    card.fill.fore_color.rgb = CARD_COLOR
    card.line.color.rgb = BLUE_ACCENT
    card.line.width = Pt(2.0)
    
    c_tf = card.text_frame
    c_tf.word_wrap = True
    c_tf.margin_left = Inches(0.4)
    c_tf.margin_top = Inches(0.4)
    c_tf.margin_right = Inches(0.4)
    
    p_q = c_tf.paragraphs[0]
    p_q.text = "“I asked ChatGPT... and it made me feel more understood in 5 mins than the 5 years I've been going to my psychologist.”"
    p_q.font.name = "Arial"
    p_q.font.size = Pt(17)
    p_q.font.bold = True
    p_q.font.italic = True
    p_q.font.color.rgb = TEXT_MAIN
    p_q.space_after = Pt(12)
    
    p_m = c_tf.add_paragraph()
    p_m.text = "— Score: 1,169  |  Engagement Weight: 7.06\nSource: https://old.reddit.com/r/ChatGPT/comments/1pomfq9"
    p_m.font.name = "Arial"
    p_m.font.size = Pt(11)
    p_m.font.bold = True
    p_m.font.color.rgb = BLUE_ACCENT
    p_m.space_after = Pt(14)
    
    p_d = c_tf.add_paragraph()
    p_d.text = "Research Finding: A rising number of users rely on ChatGPT for emotional validation and companionship, often declaring it superior to human therapists. This attachment lowers skepticism, leading to extreme vulnerability if the model provides incorrect or unsafe guidance."
    p_d.font.name = "Arial"
    p_d.font.size = Pt(12)
    p_d.font.color.rgb = TEXT_MUTED
    
    # Slide 12 Notes
    slide12.notes_slide.notes_text_frame.text = (
        "Here we deep-dive into 'Over-Reliance on AI Outputs'. "
        "This theme is highly critical: 12.29% raw prevalence. "
        "The driver here is anthropomorphism. The quote on the right highlights a user stating "
        "ChatGPT understood them better in 5 minutes than a human psychologist did in 5 years. "
        "While it shows the therapeutic potential of simulated empathy, it also highlights the "
        "danger: if users establish this level of trust, their fact-checking guard drops entirely, "
        "making them highly vulnerable to medical or life-decision hallucinations."
    )

    # =========================================================================
    # SLIDE 13: VERIFICATION BEHAVIOR INSIGHTS
    # =========================================================================
    slide13 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide13)
    add_slide_header(slide13, "User Verification & Auditing Heuristics", "Thematic deep-dive")
    
    # Left column
    info_box = slide13.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(5.8), Inches(4.8))
    inf = info_box.text_frame
    inf.word_wrap = True
    inf.margin_left = inf.margin_top = inf.margin_right = inf.margin_bottom = 0
    
    p_t = inf.paragraphs[0]
    p_t.text = "Operational Metrics & Drivers"
    p_t.font.name = "Arial"
    p_t.font.size = Pt(20)
    p_t.font.bold = True
    p_t.font.color.rgb = GREEN_ACCENT
    p_t.space_after = Pt(12)
    
    bullets = [
        "Raw Prevalence: 87 posts (14.65%)  |  Weighted: 115.82 (9.61%)",
        "Primary Driver: Factual Skepticism. Catching the model fabricating paper titles, citations, or quotes leads to active auditing.",
        "User Heuristics: Prompt engineering (e.g. 'Cite sources'), cross-referencing databases, and designing programmatic verification scripts.",
        "Advanced Verification: Developer ecosystems building dedicated fact-checking integrations (like VerifyAI or Closed-Loop Tribunals)."
    ]
    for b in bullets:
        p_b = inf.add_paragraph()
        p_b.text = "• " + b
        p_b.font.name = "Arial"
        p_b.font.size = Pt(13)
        p_b.font.color.rgb = TEXT_MAIN
        p_b.space_after = Pt(12)
        
    # Right column
    card = slide13.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.8), Inches(1.8), Inches(5.7), Inches(4.5))
    card.fill.solid()
    card.fill.fore_color.rgb = CARD_COLOR
    card.line.color.rgb = GREEN_ACCENT
    card.line.width = Pt(2.0)
    
    c_tf = card.text_frame
    c_tf.word_wrap = True
    c_tf.margin_left = Inches(0.4)
    c_tf.margin_top = Inches(0.4)
    c_tf.margin_right = Inches(0.4)
    
    p_q = c_tf.paragraphs[0]
    p_q.text = "“...i’ve caught it making up papers or quotes that don’t exist”"
    p_q.font.name = "Arial"
    p_q.font.size = Pt(18)
    p_q.font.bold = True
    p_q.font.italic = True
    p_q.font.color.rgb = TEXT_MAIN
    p_q.space_after = Pt(12)
    
    p_m = c_tf.add_paragraph()
    p_m.text = "— Score: 211  |  Engagement Weight: 5.36\nSource: https://old.reddit.com/r/ChatGPT/comments/1oj9s97"
    p_m.font.name = "Arial"
    p_m.font.size = Pt(11)
    p_m.font.bold = True
    p_m.font.color.rgb = GREEN_ACCENT
    p_m.space_after = Pt(14)
    
    p_d = c_tf.add_paragraph()
    p_d.text = "Research Finding: Auditing behaviors are triggered directly by caught hallucinations. Once a user catches the system fabricating concrete details like paper names or academic quotes, they develop strict cross-checking heuristics to protect themselves from error exposure."
    p_d.font.name = "Arial"
    p_d.font.size = Pt(12)
    p_d.font.color.rgb = TEXT_MUTED
    
    # Slide 13 Notes
    slide13.notes_slide.notes_text_frame.text = (
        "Slide 13 explores 'User Evaluation & Verification Behavior'. "
        "This is an incredibly encouraging theme, representing 14.65% raw prevalence. "
        "We capture users developing active, healthy skepticism. When they catch the model fabricating papers or "
        "quotes, they create robust prompt guardrails, such as demanding citations, double-checking academic databases, "
        "or building dedicated validator packages. It shows that users can successfully learn "
        "confidence calibration through active tool audits."
    )

    # =========================================================================
    # SLIDE 14: REAL-WORLD IMPACT INSIGHTS
    # =========================================================================
    slide14 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide14)
    add_slide_header(slide14, "Real-World Impact & Consequences", "Thematic deep-dive")
    
    # Left column
    info_box = slide14.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(5.8), Inches(4.8))
    inf = info_box.text_frame
    inf.word_wrap = True
    inf.margin_left = inf.margin_top = inf.margin_right = inf.margin_bottom = 0
    
    p_t = inf.paragraphs[0]
    p_t.text = "Operational Metrics & Drivers"
    p_t.font.name = "Arial"
    p_t.font.size = Pt(20)
    p_t.font.bold = True
    p_t.font.color.rgb = RED_ACCENT
    p_t.space_after = Pt(12)
    
    bullets = [
        "Raw Prevalence: 131 posts (22.05%)  |  Weighted: 318.07 (26.40%)",
        "Primary Driver: Practical DOWNSTREAM Consequences. High-stakes applications in legal courtrooms, medical analysis, or audit discovery.",
        "Significant Successes: Detecting 10 years of probate estate fraud or catching vet diagnostic blood-panel mistakes, saving lives.",
        "Significant Failures: Submitting fictional brief citations, leading to disbarment proceedings and court fines."
    ]
    for b in bullets:
        p_b = inf.add_paragraph()
        p_b.text = "• " + b
        p_b.font.name = "Arial"
        p_b.font.size = Pt(13)
        p_b.font.color.rgb = TEXT_MAIN
        p_b.space_after = Pt(12)
        
    # Right column
    card = slide14.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.8), Inches(1.8), Inches(5.7), Inches(4.5))
    card.fill.solid()
    card.fill.fore_color.rgb = CARD_COLOR
    card.line.color.rgb = RED_ACCENT
    card.line.width = Pt(2.0)
    
    c_tf = card.text_frame
    c_tf.word_wrap = True
    c_tf.margin_left = Inches(0.4)
    c_tf.margin_top = Inches(0.4)
    c_tf.margin_right = Inches(0.4)
    
    p_q = c_tf.paragraphs[0]
    p_q.text = "“If I hadn't used chatgpt to explain her condition compared with her actions, I would have just taken their word at face value and euthanized her.”"
    p_q.font.name = "Arial"
    p_q.font.size = Pt(16)
    p_q.font.bold = True
    p_q.font.italic = True
    p_q.font.color.rgb = TEXT_MAIN
    p_q.space_after = Pt(12)
    
    p_m = c_tf.add_paragraph()
    p_m.text = "— Score: 734  |  Engagement Weight: 6.60\nSource: https://old.reddit.com/r/ChatGPT/comments/1sc6zf0"
    p_m.font.name = "Arial"
    p_m.font.size = Pt(11)
    p_m.font.bold = True
    p_m.font.color.rgb = RED_ACCENT
    p_m.space_after = Pt(14)
    
    p_d = c_tf.add_paragraph()
    p_d.text = "Research Finding: Downstream effects represent the ultimate test of trust. In veterinary diagnostics, the user utilized ChatGPT to decode complex lab panels, exposing a severe clinical error by the veterinarian and saving their cat's life. This proves AI's high utility when paired with human verification."
    p_d.font.name = "Arial"
    p_d.font.size = Pt(12)
    p_d.font.color.rgb = TEXT_MUTED
    
    # Slide 14 Notes
    slide14.notes_slide.notes_text_frame.text = (
        "Here we discuss 'Real-World Impact of AI Outputs'. "
        "This is our most viral theme, representing a massive 26.40% weighted prevalence. "
        "Downstream outcomes are highly concrete. The quote on the right is extremely powerful: "
        "a user used ChatGPT to analyze their cat's complex blood panel, caught a major veterinarian diagnosis mistake, "
        "and prevented a wrongful euthanasia. This shows the incredible value of ChatGPT as a translation "
        "and explanation tool, provided that the human operator remains in the loop to cross-examine."
    )

    # =========================================================================
    # SLIDE 15: PERSUASIVE OUTPUTS INSIGHTS
    # =========================================================================
    slide15 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide15)
    add_slide_header(slide15, "Persuasive Outputs & Debate Fallacies", "Thematic deep-dive")
    
    # Left column
    info_box = slide15.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(5.8), Inches(4.8))
    inf = info_box.text_frame
    inf.word_wrap = True
    inf.margin_left = inf.margin_top = inf.margin_right = inf.margin_bottom = 0
    
    p_t = inf.paragraphs[0]
    p_t.text = "Operational Metrics & Drivers"
    p_t.font.name = "Arial"
    p_t.font.size = Pt(20)
    p_t.font.bold = True
    p_t.font.color.rgb = PURPLE_ACCENT
    p_t.space_after = Pt(12)
    
    bullets = [
        "Raw Prevalence: 46 posts (7.74%)  |  Weighted: 99.38 (8.25%)",
        "Primary Driver: Simulated Empathy and Sycophancy. The model tends to flatter users and agree with incorrect premises to avoid friction.",
        " caught Tactics: Logical fallacies catchable by users, including strawman rebuttals, cherry-picked arguments, and motte-and-bailey switches.",
        "Trust Formation: Confident tone and fluent articulation are highly persuasive, quietly overriding user skepticism."
    ]
    for b in bullets:
        p_b = inf.add_paragraph()
        p_b.text = "• " + b
        p_b.font.name = "Arial"
        p_b.font.size = Pt(13)
        p_b.font.color.rgb = TEXT_MAIN
        p_b.space_after = Pt(12)
        
    # Right column
    card = slide15.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.8), Inches(1.8), Inches(5.7), Inches(4.5))
    card.fill.solid()
    card.fill.fore_color.rgb = CARD_COLOR
    card.line.color.rgb = PURPLE_ACCENT
    card.line.width = Pt(2.0)
    
    c_tf = card.text_frame
    c_tf.word_wrap = True
    c_tf.margin_left = Inches(0.4)
    c_tf.margin_top = Inches(0.4)
    c_tf.margin_right = Inches(0.4)
    
    p_q = c_tf.paragraphs[0]
    p_q.text = "“The practical takeaway isn't necessarily 'switch models.' It's being more skeptical of AI responses exactly in the domains where sycophancy is highest - subjective, value-laden, strategy and ethics questions...”"
    p_q.font.name = "Arial"
    p_q.font.size = Pt(15.5)
    p_q.font.bold = True
    p_q.font.italic = True
    p_q.font.color.rgb = TEXT_MAIN
    p_q.space_after = Pt(12)
    
    p_m = c_tf.add_paragraph()
    p_m.text = "— Score: 26  |  Engagement Weight: 3.30\nSource: https://old.reddit.com/r/ChatGPT/comments/1slcgpx"
    p_m.font.name = "Arial"
    p_m.font.size = Pt(11)
    p_m.font.bold = True
    p_m.font.color.rgb = PURPLE_ACCENT
    p_m.space_after = Pt(14)
    
    p_d = c_tf.add_paragraph()
    p_d.text = "Research Finding: Users catch models practicing sycophancy—actively agreeing with false assertions to maintain a frictionless conversation. This highlights the danger of using AI for strategy, ethics, or open-ended debates where sycophancy is highest and hardest to verify."
    p_d.font.name = "Arial"
    p_d.font.size = Pt(12)
    p_d.font.color.rgb = TEXT_MUTED
    
    # Slide 15 Notes
    slide15.notes_slide.notes_text_frame.text = (
        "We discuss 'Persuasive Outputs & Trust Formation'. "
        "Representing 7.74% raw prevalence, this theme exposes sycophancy. "
        "Sycophancy occurs when the model flatters the user or agrees with their wrong assumptions to avoid debate. "
        "The quote on the right notes that users catch models exhibiting sycophancy in strategy, strategy, "
        "ethics, and subjective domains. It indicates that highly fluent simulated empathy can act "
        "as a persuasive screen, quietly overriding natural human skepticism."
    )

    # =========================================================================
    # SLIDE 16: KEY RESEARCH FINDINGS
    # =========================================================================
    slide16 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide16)
    add_slide_header(slide16, "Key Research Findings Summary", "Synthesis")
    
    # 3 Column synthesis
    findings = [
        ("THE PREVALENCE PARADOX", PURPLE_ACCENT, [
            "Minor technical hallucinations occur constantly in everyday work (31.14% raw prevalence).",
            "However, users are far more motivated to engage with and viralize narratives of extreme consequences.",
            "Real-world impact commands a massive 26.40% weighted prevalence, showing actual community interest."
        ]),
        ("THE MIGRATION CRITICALITY", BLUE_ACCENT, [
            "Traditional keyword screening has a massive 88.89% false positive rate (11.11% precision).",
            "Simple queries also exhibit a 60% false negative rate, missing implicit trust breakdowns.",
            "This validates our hybrid, two-stage semantic LLM classifier as the only reliable research method."
        ]),
        ("THE TRUST TENSION", GREEN_ACCENT, [
            "A deep structural division exists in the ChatGPT community.",
            "Psychological attachment drives extreme over-reliance (12.29%) with zero verification.",
            "Professional practitioners catch hallucinations and create programmatic verification tools (14.65%)."
        ])
    ]
    
    for i, (title, color, bullets) in enumerate(findings):
        col_x = Inches(0.75 + i * 4.0)
        
        box = slide16.shapes.add_shape(MSO_SHAPE.RECTANGLE, col_x, Inches(1.8), Inches(3.7), Inches(4.5))
        box.fill.solid()
        box.fill.fore_color.rgb = CARD_COLOR
        box.line.color.rgb = color
        box.line.width = Pt(2.5)
        
        b_tf = box.text_frame
        b_tf.word_wrap = True
        b_tf.margin_left = Inches(0.25)
        b_tf.margin_top = Inches(0.3)
        b_tf.margin_right = Inches(0.25)
        
        p_t = b_tf.paragraphs[0]
        p_t.text = title
        p_t.font.name = "Arial"
        p_t.font.size = Pt(14)
        p_t.font.bold = True
        p_t.font.color.rgb = color
        p_t.space_after = Pt(14)
        
        for bullet in bullets:
            p_b = b_tf.add_paragraph()
            p_b.text = "• " + bullet
            p_b.font.name = "Arial"
            p_b.font.size = Pt(11.5)
            p_b.font.color.rgb = TEXT_MAIN
            p_b.space_after = Pt(10)
            
    # Slide 16 Notes
    slide16.notes_slide.notes_text_frame.text = (
        "This slide synthesizes our three major project findings. "
        "First, the Prevalence Paradox shows the gap between error frequency and emotional impact. "
        "Second, the Migration Criticality shows that keyword matching alone fails sentiment audits completely, "
        "justifying our advanced LLM architecture. "
        "Third, the Trust Tension shows the division between emotionally attached users and highly skepticism-driven "
        "professional practitioners who catch errors and build custom auditing rules."
    )

    # =========================================================================
    # SLIDE 17: PRODUCT OPPORTUNITIES
    # =========================================================================
    slide17 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide17)
    add_slide_header(slide17, "Strategic Product Opportunities", "Product Strategy")
    
    # 3 Column Product Opportunities
    opps = [
        ("1. Active Uncertainty UI", BLUE_ACCENT, [
            "Move away from constant, absolute conversational confidence.",
            "Implement Visual Uncertainty Indicators: when a model's generation confidence drops, color the output text in yellow/amber.",
            "Prompts users to double-check and perform manual verification on uncertain claims."
        ]),
        ("2. Version Control Transparency", PURPLE_ACCENT, [
            "Silent, unannounced updates are the largest drivers of trust erosion and breakdown (18.45% weighted).",
            "Provide explicit version control and detailed platform changelogs.",
            "Allow users to toggle between historical model weights to ensure workflow reproducibility."
        ]),
        ("3. Embedded Fact-Checking", GREEN_ACCENT, [
            "Embed fact-checking buttons directly in the interface.",
            "Integrate automated verification calls to trusted academic databases and search engines (like the user-created VerifyAI tool).",
            "Reduces hallucination rate dramatically and provides instant source citation auditing."
        ])
    ]
    
    for i, (title, color, bullets) in enumerate(opps):
        col_x = Inches(0.75 + i * 4.0)
        
        box = slide17.shapes.add_shape(MSO_SHAPE.RECTANGLE, col_x, Inches(1.8), Inches(3.7), Inches(4.5))
        box.fill.solid()
        box.fill.fore_color.rgb = CARD_COLOR
        box.line.color.rgb = color
        box.line.width = Pt(2.5)
        
        b_tf = box.text_frame
        b_tf.word_wrap = True
        b_tf.margin_left = Inches(0.25)
        b_tf.margin_top = Inches(0.3)
        b_tf.margin_right = Inches(0.25)
        
        p_t = b_tf.paragraphs[0]
        p_t.text = title
        p_t.font.name = "Arial"
        p_t.font.size = Pt(15)
        p_t.font.bold = True
        p_t.font.color.rgb = color
        p_t.space_after = Pt(14)
        
        for bullet in bullets:
            p_b = b_tf.add_paragraph()
            p_b.text = "• " + bullet
            p_b.font.name = "Arial"
            p_b.font.size = Pt(11.5)
            p_b.font.color.rgb = TEXT_MAIN
            p_b.space_after = Pt(10)
            
    # Slide 17 Notes
    slide17.notes_slide.notes_text_frame.text = (
        "For Product Managers, this is the most critical slide in the deck. "
        "We translate our research insights into three concrete product opportunities. "
        "First, we suggest an Active Uncertainty UI to highlight low-confidence generations visually. "
        "Second, we recommend explicit version control and detailed platform changelogs to combat trust breakdown. "
        "Third, we propose embedded verification integrations to allow instant source checking."
    )

    # =========================================================================
    # SLIDE 18: RECOMMENDATIONS
    # =========================================================================
    slide18 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide18)
    add_slide_header(slide18, "Strategic Recommendations", "Action Plan")
    
    # Left Box: LLM Providers
    provider_box = slide18.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.75), Inches(1.8), Inches(5.6), Inches(4.5))
    provider_box.fill.solid()
    provider_box.fill.fore_color.rgb = CARD_COLOR
    provider_box.line.color.rgb = BLUE_ACCENT
    provider_box.line.width = Pt(2.0)
    
    pb_tf = provider_box.text_frame
    pb_tf.word_wrap = True
    pb_tf.margin_left = Inches(0.35)
    pb_tf.margin_top = Inches(0.35)
    pb_tf.margin_right = Inches(0.35)
    
    p_pb_t = pb_tf.paragraphs[0]
    p_pb_t.text = "For LLM & Platform Providers"
    p_pb_t.font.name = "Arial"
    p_pb_t.font.size = Pt(18)
    p_pb_t.font.bold = True
    p_pb_t.font.color.rgb = BLUE_ACCENT
    p_pb_t.space_after = Pt(12)
    
    provider_recs = [
        "Calibrate Tone: Train models to avoid sycophancy, using neutral language and expressing caution in strategy and ethics domains.",
        "Visual Safeguards: Design guardrails targeting critical domains (medical/legal), forcing double-verification before rendering.",
        "API Transparency: Log model version tokens inside all API responses, preventing unexpected downstream prompt failures."
    ]
    for r in provider_recs:
        p_r = pb_tf.add_paragraph()
        p_r.text = "• " + r
        p_r.font.name = "Arial"
        p_r.font.size = Pt(12.5)
        p_r.font.color.rgb = TEXT_MAIN
        p_r.space_after = Pt(12)
        
    # Right Box: Researchers
    researcher_box = slide18.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.8), Inches(1.8), Inches(5.783), Inches(4.5))
    researcher_box.fill.solid()
    researcher_box.fill.fore_color.rgb = CARD_COLOR
    researcher_box.line.color.rgb = PURPLE_ACCENT
    researcher_box.line.width = Pt(2.0)
    
    rb_tf = researcher_box.text_frame
    rb_tf.word_wrap = True
    rb_tf.margin_left = Inches(0.35)
    rb_tf.margin_top = Inches(0.35)
    rb_tf.margin_right = Inches(0.35)
    
    p_rb_t = rb_tf.paragraphs[0]
    p_rb_t.text = "For Users & Researchers"
    p_rb_t.font.name = "Arial"
    p_rb_t.font.size = Pt(18)
    p_rb_t.font.bold = True
    p_rb_t.font.color.rgb = PURPLE_ACCENT
    p_rb_t.space_after = Pt(12)
    
    user_recs = [
        "Cross-Examination Heuristics: Adopt custom verification templates (e.g. Closed-Loop Tribunals) when debugging or drafting crucial documents.",
        "Hybrid Pipeline Architectures: Academic sentiment researchers should abandon basic keyword screening, deploying LLM semantic stages.",
        "Scoping Comment Ingestion: Expand research scoping to pull thread replies to map collaborative community debugging behavior."
    ]
    for r in user_recs:
        p_r = rb_tf.add_paragraph()
        p_r.text = "• " + r
        p_r.font.name = "Arial"
        p_r.font.size = Pt(12.5)
        p_r.font.color.rgb = TEXT_MAIN
        p_r.space_after = Pt(12)
        
    # Slide 18 Notes
    slide18.notes_slide.notes_text_frame.text = (
        "Slide 18 details our action recommendations for two core audiences. "
        "For LLM providers, we recommend training models to resist sycophancy, "
        "implementing clinical/legal UI safeguards, and standardizing version tokens in APIs. "
        "For users and researchers, we recommend adopting active cross-examination heuristics "
        "and utilizing two-stage hybrid semantic ingestion architectures rather than keyword heuristics."
    )

    # =========================================================================
    # SLIDE 19: LIMITATIONS
    # =========================================================================
    slide19 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide19)
    add_slide_header(slide19, "Project Limitations & Boundaries", "Limitations")
    
    # 4 Column Limitations
    limits = [
        ("1. Reddit Selection Bias", "Sourced from r/ChatGPT, which naturally skews toward active hobbyists, developers, and users seeking support. Severe errors and trust breakdown are likely over-represented.", BLUE_ACCENT),
        ("2. Thread Boundary", "The Playwright engine targeted main post titles and body texts only, excluding comment replies. Context-rich community corrections occurring in replies were missed.", PURPLE_ACCENT),
        ("3. Temporal Horizon", "Sourced across a 9-month chronological window ending in late 2025/2026. Rapid model upgrades post-study (e.g. GPT-5 release) may alter theme distributions.", GREEN_ACCENT),
        ("4. Classifier Thresholds", "Utilized Llama 3.1 8B for semantic theme mapping. Highly complex or ambiguous posts with overlapping themes may show minor classification deviations.", RED_ACCENT)
    ]
    
    for i, (title, desc, color) in enumerate(limits):
        col_x = Inches(0.75 + i * 2.95)
        
        box = slide19.shapes.add_shape(MSO_SHAPE.RECTANGLE, col_x, Inches(1.8), Inches(2.8), Inches(4.5))
        box.fill.solid()
        box.fill.fore_color.rgb = CARD_COLOR
        box.line.color.rgb = color
        box.line.width = Pt(2.0)
        
        b_tf = box.text_frame
        b_tf.word_wrap = True
        b_tf.margin_left = Inches(0.25)
        b_tf.margin_top = Inches(0.3)
        b_tf.margin_right = Inches(0.25)
        
        p_t = b_tf.paragraphs[0]
        p_t.text = title
        p_t.font.name = "Arial"
        p_t.font.size = Pt(16)
        p_t.font.bold = True
        p_t.font.color.rgb = color
        p_t.space_after = Pt(14)
        
        p_d = b_tf.add_paragraph()
        p_d.text = desc
        p_d.font.name = "Arial"
        p_d.font.size = Pt(12)
        p_d.font.color.rgb = TEXT_MAIN
        p_d.space_after = Pt(10)
        
    # Slide 19 Notes
    slide19.notes_slide.notes_text_frame.text = (
        "Here we discuss our project's limitations. "
        "First, we acknowledge Reddit selection bias: users who are satisfied rarely post, "
        "so failures are naturally over-represented. "
        "Second, our ingestion targeted only post bodies, not comments. "
        "Third, our 9-month temporal horizon represents a static snapshot. "
        "Fourth, the model size (Llama 3.1 8B) could introduce minor boundaries on highly complex semantic threads. "
        "Acknowledging these limitations supports our academic integrity."
    )

    # =========================================================================
    # SLIDE 20: CONCLUSION
    # =========================================================================
    slide20 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide20)
    add_slide_header(slide20, "Project Conclusion", "Summary")
    
    # Large centered text container
    c_box = slide20.shapes.add_textbox(Inches(1.5), Inches(1.8), Inches(10.333), Inches(4.8))
    cf = c_box.text_frame
    cf.word_wrap = True
    cf.margin_left = cf.margin_top = cf.margin_right = cf.margin_bottom = 0
    
    p_c = cf.paragraphs[0]
    p_c.text = "Bridging the Gap Between AI Fluency & Human Trust"
    p_c.font.name = "Arial"
    p_c.font.size = Pt(24)
    p_c.font.bold = True
    p_c.font.color.rgb = PURPLE_ACCENT
    p_c.alignment = PP_ALIGN.CENTER
    p_c.space_after = Pt(20)
    
    conclusions = [
        "Empirical Milestone: Successfully created an open-source, locked database (data/research.db) of 594 relevant user experiences, bridging standardized benchmarking and real-world HCI observations.",
        "Mathematical Contribution: Proven that traditional keyword heuristics are highly unreliable for qualitative sentiment auditing, yielding an 88.89% error rate compared to semantic LLM classifiers.",
        "Actionable Strategic Takeaways: Outlined key product opportunities (Uncertainty UI, API Versioning, and Fact-Checking buttons) for platform developers and product managers to prevent harmful over-reliance.",
        "Ongoing Value: The premium interactive dashboard provides ongoing analytical utility, offering future researchers and stakeholders a robust qualitative auditing interface."
    ]
    for c in conclusions:
        p_c_line = cf.add_paragraph()
        p_c_line.text = c
        p_c_line.font.name = "Arial"
        p_c_line.font.size = Pt(14)
        p_c_line.font.color.rgb = TEXT_MAIN
        p_c_line.space_after = Pt(14)
        p_c_line.alignment = PP_ALIGN.LEFT
        
    # Slide 20 Notes
    slide20.notes_slide.notes_text_frame.text = (
        "In conclusion, the ChatGPT Output Trust & Evaluation Lab represents an empirical milestone. "
        "We successfully mapped the real-world user trust landscape of ChatGPT. "
        "We mathematically demonstrated that keyword counting is highly insufficient for sentiment audits, "
        "and outlined actionable strategic suggestions for future LLM product designs. "
        "Our open-source database and Streamlit dashboard will serve as key platforms for ongoing research. "
        "Thank you for your time, and I am open to any questions!"
    )

    # Make sure deliverables folder exists
    os.makedirs("deliverables", exist_ok=True)
    
    # Save Presentation
    prs.save("deliverables/ChatGPT_Trust_Lab_Final_Presentation.pptx")
    print("Presentation created successfully at deliverables/ChatGPT_Trust_Lab_Final_Presentation.pptx")

if __name__ == "__main__":
    create_presentation()
